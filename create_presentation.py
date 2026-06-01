"""
PowerPoint Presentation Generator — Cloud Computing Web Scraper Project

Creates a presentation about the web scraper project.
Uses a user-provided .pptx template if given, otherwise creates from scratch.

Usage:
    python create_presentation.py
    python create_presentation.py --template my_template.pptx
    python create_presentation.py --template my_template.pptx --output presentation.pptx

Requirements:
    pip install python-pptx
"""

import argparse
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("ERROR: python-pptx is required. Install with:")
    print("  pip install python-pptx")
    sys.exit(1)


# ─── Slide Content ────────────────────────────────────────────────

SLIDES = [
    {
        "title": "Cloud Computing Web Scraper",
        "subtitle": "A Serverless Mastodon Data Pipeline on Google Cloud Platform",
        "layout": "title",
    },
    {
        "title": "Project Overview",
        "content": [
            "Modular Python web scraper with Mastodon API support",
            "Deployed across 5 distinct GCP resources (no VMs)",
            "Fully serverless, event-driven architecture",
            "Scrapes real-time public data from the Fediverse",
            "Exports structured data as JSON to Cloud Storage",
        ],
        "layout": "content",
    },
    {
        "title": "Architecture Diagram",
        "content": [
            "Cloud Scheduler (cron every 10 min)",
            "    \u2193 publishes message",
            "Pub/Sub Topic (scraper-trigger)",
            "    \u2193 triggers function",
            "Cloud Function (mastodon-scraper)",
            "    \u2193 fetches data from Mastodon API",
            "    \u2193 stores result",
            "Cloud Storage (JSON output bucket)",
            "",
            "Cloud Logging captures all execution logs automatically",
        ],
        "layout": "content",
    },
    {
        "title": "5 GCP Resources Used",
        "content": [
            "1. Cloud Functions \u2014 Serverless compute running the Python scraper",
            "2. Cloud Storage \u2014 Object storage for scraped JSON data",
            "3. Pub/Sub \u2014 Event-driven messaging between services",
            "4. Cloud Scheduler \u2014 Managed cron job triggering the pipeline",
            "5. Cloud Logging \u2014 Centralized log collection & monitoring",
        ],
        "layout": "content",
    },
    {
        "title": "Resource 1: Cloud Functions",
        "content": [
            "Serverless compute \u2014 no VM provisioning needed",
            "Runtime: Python 3.11",
            "Trigger: Pub/Sub topic subscription",
            "Memory: 256 MB | Timeout: 60 seconds",
            "Entry point: scrape_mastodon(event, context)",
            "Scrapes Mastodon API and uploads result to GCS",
            "Scales to zero when not in use (cost-efficient)",
        ],
        "layout": "content",
    },
    {
        "title": "Resource 2: Cloud Storage",
        "content": [
            "Object storage for pipeline output",
            "Bucket: <project-id>-scraper-output",
            "Path structure: scrapes/mastodon_<action>_<timestamp>.json",
            "Stores structured JSON with metadata",
            "Durable, highly available, low-cost",
            "Can be queried with BigQuery for analytics",
        ],
        "layout": "content",
    },
    {
        "title": "Resource 3: Pub/Sub",
        "content": [
            "Asynchronous messaging service",
            "Topic: scraper-trigger",
            "Decouples the scheduler from the function",
            "Enables reliable at-least-once delivery",
            "Allows multiple subscribers if pipeline grows",
            "Handles backpressure and retries automatically",
        ],
        "layout": "content",
    },
    {
        "title": "Resource 4: Cloud Scheduler",
        "content": [
            "Fully managed cron job service",
            "Schedule: */10 * * * * (every 10 minutes)",
            "Publishes a message to Pub/Sub on each tick",
            "No server to maintain \u2014 just a schedule definition",
            "Can be paused/resumed from the console",
            "Supports HTTP, Pub/Sub, and App Engine targets",
        ],
        "layout": "content",
    },
    {
        "title": "Resource 5: Cloud Logging",
        "content": [
            "Automatic log capture from Cloud Functions",
            "All print() statements appear in Logs Explorer",
            "Structured logs with timestamps and severity",
            "Enables debugging and monitoring",
            "Can set up alerts on errors",
            "Retained for 30 days by default",
        ],
        "layout": "content",
    },
    {
        "title": "The Scraper: Mastodon API",
        "content": [
            "Mastodon: open-source, decentralized social network",
            "Public API \u2014 no authentication required for public data",
            "Endpoints used:",
            "  \u2022 GET /api/v1/trends/statuses (trending posts)",
            "  \u2022 GET /api/v1/timelines/tag/:hashtag",
            "  \u2022 GET /api/v1/accounts/lookup",
            "  \u2022 GET /api/v2/instance (server metadata)",
            "Respects rate limits (1 req/sec delay)",
            "Returns JSON with text, media URLs, metadata",
        ],
        "layout": "content",
    },
    {
        "title": "Data Flow Example",
        "content": [
            "1. Scheduler fires at :00, :10, :20, :30, :40, :50",
            "2. Pub/Sub receives message \u2192 triggers Cloud Function",
            "3. Function calls: GET mastodon.social/api/v1/trends/statuses?limit=1",
            "4. Parses response: extracts text, author, media, metrics",
            "5. Wraps in metadata envelope (timestamp, instance, action)",
            "6. Uploads JSON to gs://bucket/scrapes/mastodon_trends_20260601_120000.json",
            "7. Logs success to Cloud Logging",
        ],
        "layout": "content",
    },
    {
        "title": "Technology Stack",
        "content": [
            "Language: Python 3.11",
            "HTTP Client: requests library",
            "Cloud SDK: google-cloud-storage",
            "Deployment: gcloud CLI",
            "Source Control: GitHub",
            "Local UI: Tkinter (glassmorphism theme)",
            "Export Formats: JSON, TXT, HTML, Binary",
        ],
        "layout": "content",
    },
    {
        "title": "Deployment Process",
        "content": [
            "1. Code pushed to GitHub repository",
            "2. Clone to Google Cloud Shell",
            "3. Run deploy.sh \u2014 single script sets up all 5 resources:",
            "   \u2022 Enables required APIs",
            "   \u2022 Creates GCS bucket",
            "   \u2022 Creates Pub/Sub topic",
            "   \u2022 Deploys Cloud Function",
            "   \u2022 Creates Cloud Scheduler job",
            "4. Pipeline runs automatically every 10 minutes",
        ],
        "layout": "content",
    },
    {
        "title": "Live Demo Results",
        "content": [
            "Successfully scraped trending Mastodon posts",
            "Data stored in Cloud Storage bucket",
            "Logs visible in Cloud Logging",
            "Pipeline runs every 10 minutes automatically",
            "Zero VMs \u2014 fully serverless architecture",
            "Cost: essentially free at this scale (free tier)",
        ],
        "layout": "content",
    },
    {
        "title": "Summary & Conclusion",
        "content": [
            "Built a functional web scraper deployed on GCP",
            "Uses 5 distinct cloud resources (no VMs)",
            "Event-driven, serverless, scalable architecture",
            "Real-world data source (Mastodon public API)",
            "Automated pipeline with scheduled execution",
            "Infrastructure as code (deploy.sh)",
        ],
        "layout": "content",
    },
    {
        "title": "Thank You",
        "subtitle": "Questions?",
        "layout": "title",
    },
]


# ─── Presentation Builder ─────────────────────────────────────────

def create_presentation(template_path=None, output_path="presentation.pptx"):
    """Build the PowerPoint presentation."""

    if template_path and os.path.exists(template_path):
        print(f"[*] Using template: {template_path}")
        prs = Presentation(template_path)
    else:
        if template_path:
            print(f"[!] Template not found: {template_path}, creating from scratch")
        prs = Presentation()

    # Get available layouts
    layouts = prs.slide_layouts
    print(f"[*] Available slide layouts ({len(layouts)}):")
    for i, layout in enumerate(layouts):
        print(f"    [{i}] {layout.name}")

    # Try to identify useful layouts
    title_layout = None
    content_layout = None

    for i, layout in enumerate(layouts):
        name = layout.name.lower()
        if "title slide" in name or (i == 0 and title_layout is None):
            title_layout = layout
        elif "title and content" in name or "content" in name:
            content_layout = layout
        elif "title only" in name and content_layout is None:
            content_layout = layout

    # Fallbacks
    if title_layout is None:
        title_layout = layouts[0]
    if content_layout is None:
        content_layout = layouts[1] if len(layouts) > 1 else layouts[0]

    print(f"[*] Using title layout: {title_layout.name}")
    print(f"[*] Using content layout: {content_layout.name}")
    print(f"[*] Creating {len(SLIDES)} slides...")

    for slide_data in SLIDES:
        if slide_data["layout"] == "title":
            slide = prs.slides.add_slide(title_layout)
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
            # Set subtitle if placeholder exists
            subtitle_text = slide_data.get("subtitle", "")
            if subtitle_text:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # subtitle placeholder
                        shape.text = subtitle_text
                        break

        elif slide_data["layout"] == "content":
            slide = prs.slides.add_slide(content_layout)
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]

            # Set body content
            body_text = "\n".join(slide_data.get("content", []))
            body_set = False
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # body placeholder
                    tf = shape.text_frame
                    tf.clear()
                    for i, line in enumerate(slide_data.get("content", [])):
                        if i == 0:
                            tf.paragraphs[0].text = line
                            tf.paragraphs[0].font.size = Pt(18)
                        else:
                            p = tf.add_paragraph()
                            p.text = line
                            p.font.size = Pt(18)
                            p.space_before = Pt(6)
                    body_set = True
                    break

            # If no body placeholder, add a text box
            if not body_set:
                left = Inches(0.8)
                top = Inches(1.8)
                width = Inches(8.5)
                height = Inches(5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, line in enumerate(slide_data.get("content", [])):
                    if i == 0:
                        tf.paragraphs[0].text = line
                        tf.paragraphs[0].font.size = Pt(18)
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(18)
                        p.space_before = Pt(6)

    # Save
    prs.save(output_path)
    print(f"\n[*] Presentation saved to: {output_path}")
    print(f"[*] Total slides: {len(prs.slides)}")
    print(f"\n[*] Done! Open {output_path} in PowerPoint.")
    print("[*] You can also paste this into Copilot with your template for styling.")


def main():
    parser = argparse.ArgumentParser(description="Generate project presentation (PPTX)")
    parser.add_argument("--template", default=None,
                        help="Path to a .pptx template file to use as base")
    parser.add_argument("--output", default="presentation.pptx",
                        help="Output filename (default: presentation.pptx)")
    args = parser.parse_args()

    create_presentation(template_path=args.template, output_path=args.output)


if __name__ == "__main__":
    main()
